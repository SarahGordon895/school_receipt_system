from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT_DIR = ROOT / "docs" / "presentation"
ASSETS = OUT_DIR / "assets"
DIAGRAMS = ROOT / "docs" / "diagrams"
TEMPLATE = Path("/Users/apple/Desktop/IS098/UDSM_Branding_PPT.pptx")

SLIDE_W = 13.333
SLIDE_H = 7.5

GREEN = RGBColor(9, 74, 55)
GREEN_DARK = RGBColor(5, 52, 39)
GREEN_SOFT = RGBColor(225, 239, 233)
GOLD = RGBColor(205, 160, 39)
RED = RGBColor(211, 55, 66)
BLUE = RGBColor(0, 112, 192)
INK = RGBColor(28, 43, 39)
MUTED = RGBColor(91, 109, 103)
PAPER = RGBColor(247, 249, 247)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(214, 223, 218)


prs = Presentation(str(TEMPLATE))
for slide_id in list(prs.slides._sldIdLst):
    prs.part.drop_rel(slide_id.rId)
    prs.slides._sldIdLst.remove(slide_id)
BLANK = prs.slide_layouts[6]


def rect(slide, x, y, w, h, fill, line=None, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def textbox(slide, text, x, y, w, h, size=18, color=INK, bold=False,
            font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
            margin=0.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def rich_textbox(slide, runs, x, y, w, h, size=18, color=INK,
                 align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.vertical_anchor = valign
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.05)
    p = frame.paragraphs[0]
    p.alignment = align
    for item in runs:
        r = p.add_run()
        r.text = item["text"]
        r.font.name = "Aptos"
        r.font.size = Pt(item.get("size", size))
        r.font.bold = item.get("bold", False)
        r.font.color.rgb = item.get("color", color)
    return box


def bullets(slide, items, x, y, w, h, size=17, color=INK, bullet_color=None,
            spacing=8, level_indent=0.22):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    for idx, item in enumerate(items):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"•  {text}"
        p.level = level
        p.font.name = "Aptos"
        p.font.size = Pt(size - level * 2)
        p.font.color.rgb = color
        p.space_after = Pt(spacing)
        p.line_spacing = 1.05
        p.left_margin = Inches(level * level_indent)
    return box


def label(slide, text, x, y, w, fill=GREEN_SOFT, color=GREEN, size=12):
    rect(slide, x, y, w, 0.32, fill, fill, True)
    textbox(slide, text, x + 0.06, y + 0.02, w - 0.12, 0.25, size=size,
            color=color, bold=True, align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE)


def title_bar(slide, title, number, subtitle=None):
    rect(slide, 2.30, 0, 11.03, 1.47, BLUE, BLUE, False)
    textbox(slide, "UNIVERSITY OF DAR ES SALAAM", 2.48, 0.12, 4.7, 0.24,
            size=9, color=WHITE, bold=True)
    textbox(slide, title, 2.46, 0.38, 9.65, 0.50, size=22, color=WHITE, bold=True)
    textbox(slide, f"{number:02d}", 12.12, 0.24, 0.62, 0.34, size=12,
            color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
    if subtitle:
        textbox(slide, subtitle, 2.48, 0.93, 9.65, 0.28, size=9.5,
                color=RGBColor(226, 239, 248))
    rect(slide, 0, 1.46, SLIDE_W, 0.04, GOLD, GOLD, False)


def footer(slide, number, source="Final Year Project Report & implemented system · July 2026"):
    textbox(slide, source, 0.58, 7.12, 10.8, 0.22, size=8.5, color=MUTED)
    textbox(slide, str(number), 12.08, 7.09, 0.45, 0.22, size=9, color=MUTED,
            align=PP_ALIGN.RIGHT)


def add_picture_contain(slide, image_path, x, y, w, h, border=True):
    image_path = Path(image_path)
    with Image.open(image_path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    px, py = x + (w - pw) / 2, y + (h - ph) / 2
    if border:
        rect(slide, x, y, w, h, WHITE, LINE, True)
    slide.shapes.add_picture(str(image_path), Inches(px), Inches(py), Inches(pw), Inches(ph))


def add_picture_crop(slide, image_path, x, y, w, h, border=True):
    image_path = Path(image_path)
    with Image.open(image_path) as im:
        iw, ih = im.size
    target = w / h
    ratio = iw / ih
    pic = slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), Inches(w), Inches(h))
    if ratio > target:
        visible = target / ratio
        crop = (1 - visible) / 2
        pic.crop_left = crop
        pic.crop_right = crop
    elif ratio < target:
        visible = ratio / target
        crop = (1 - visible) / 2
        pic.crop_top = crop
        pic.crop_bottom = crop
    if border:
        frame = rect(slide, x, y, w, h, WHITE, LINE, True)
        slide.shapes._spTree.remove(frame._element)
        slide.shapes._spTree.insert(2, frame._element)
    return pic


def add_card(slide, title, body, x, y, w, h, accent=GREEN, number=None):
    rect(slide, x, y, w, h, WHITE, LINE, True)
    rect(slide, x, y, 0.08, h, accent, accent, False)
    if number:
        rect(slide, x + 0.20, y + 0.18, 0.44, 0.44, accent, accent, True)
        textbox(slide, str(number), x + 0.22, y + 0.20, 0.40, 0.36, size=13,
                color=WHITE, bold=True, align=PP_ALIGN.CENTER,
                valign=MSO_ANCHOR.MIDDLE)
        tx = x + 0.78
    else:
        tx = x + 0.25
    textbox(slide, title, tx, y + 0.16, w - (tx - x) - 0.2, 0.36,
            size=15, color=GREEN_DARK, bold=True)
    textbox(slide, body, x + 0.25, y + 0.62, w - 0.45, h - 0.75,
            size=11.5, color=MUTED)


def add_metric(slide, value, label_text, x, y, w, accent=GREEN):
    rect(slide, x, y, w, 1.05, WHITE, LINE, True)
    textbox(slide, value, x + 0.12, y + 0.13, w - 0.24, 0.42, size=23,
            color=accent, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, label_text, x + 0.12, y + 0.61, w - 0.24, 0.25, size=10,
            color=MUTED, align=PP_ALIGN.CENTER)


# Slide 1 — Title
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 1.47, SLIDE_W, 6.03, PAPER, PAPER, False)
rect(slide, 0, 1.47, 3.95, 6.03, BLUE, BLUE, False)
rect(slide, 0, 1.46, SLIDE_W, 0.04, GOLD, GOLD, False)
textbox(slide, "FINAL YEAR PROJECT\nPRESENTATION", 0.58, 2.02, 2.95, 1.2,
        size=22, color=WHITE, bold=True)
textbox(slide, "College of Information and Communication Technologies\nDepartment of Computer Science and Engineering",
        0.58, 5.70, 2.95, 0.85, size=11.0, color=RGBColor(226, 239, 248))
textbox(slide, "July 2026", 0.58, 6.75, 2.0, 0.3, size=11, color=GOLD, bold=True)
textbox(slide, "WEB-BASED FEE TRACKING\nAND REMINDER SYSTEM",
        4.52, 2.00, 7.80, 1.35, size=30, color=BLUE, bold=True)
textbox(slide, "A focused fee management, receipt, reporting and parent communication platform for Mbonea Secondary School",
        4.55, 3.46, 7.45, 0.8, size=16, color=MUTED)
rect(slide, 4.55, 4.26, 7.55, 0.02, GOLD, GOLD, False)
textbox(slide, "Presented by", 4.55, 4.50, 1.4, 0.3, size=12, color=MUTED, bold=True)
textbox(slide,
        "Innocent Richard Mkumbo  ·  2024-02-00404\n"
        "Sarah George Gordon  ·  2024-02-00133\n"
        "Charles Dani Chaula  ·  2024-02-00080",
        4.55, 4.86, 5.55, 1.22, size=15.5, color=INK)
textbox(slide, "Supervisor", 10.20, 4.50, 1.4, 0.3, size=12, color=MUTED, bold=True)
textbox(slide, "Sir Paul Haule", 10.20, 4.88, 1.8, 0.4, size=16, color=BLUE, bold=True)
label(slide, "Diploma in Computer Science", 9.35, 6.26, 2.75,
      RGBColor(226, 239, 248), BLUE, 11)


# Slide 2 — Background
slide = prs.slides.add_slide(BLANK)
title_bar(slide, "Project background", 2, "Why Mbonea Secondary School needed a focused fee information system")
textbox(slide, "The existing environment", 0.62, 1.52, 4.2, 0.36, size=18, color=GREEN_DARK, bold=True)
bullets(slide, [
    "Fee records were spread across receipt books, notebooks and separate spreadsheets.",
    "Staff repeatedly searched and reconciled records to establish a student's current balance.",
    "Parents often depended on students, phone calls or school visits for payment information.",
    "As transactions increased, inconsistent balances and delayed communication became more likely.",
], 0.66, 1.98, 5.45, 3.6, size=16, spacing=12)
rect(slide, 0.62, 5.72, 5.55, 0.92, GREEN_SOFT, GREEN_SOFT, True)
rich_textbox(slide, [
    {"text": "Project response: ", "bold": True, "color": GREEN_DARK},
    {"text": "one web application for students, fees, payments, receipts, reports and parent communication.", "color": INK},
], 0.84, 5.94, 5.12, 0.5, size=14)

steps = [
    ("Paper records", "Slow search"),
    ("Separate files", "Different balances"),
    ("Informal messages", "Delayed parent updates"),
    ("Central system", "Reliable information"),
]
for i, (title, body) in enumerate(steps):
    y = 1.65 + i * 1.25
    accent = GREEN if i == 3 else GOLD
    add_card(slide, title, body, 6.65, y, 5.65, 0.92, accent=accent, number=i + 1)
footer(slide, 2)


# Slide 3 — Problem
slide = prs.slides.add_slide(BLANK)
title_bar(slide, "Problem statement", 3, "Fragmented records affected accuracy, workload, reporting and parent trust")
cards = [
    ("No single source of truth", "Payment information could exist in several places, leaving balances out of date.", RED),
    ("Manual calculation workload", "Bursars spent time searching, combining and recalculating fee records.", GOLD),
    ("Slow parent communication", "Deadlines, confirmations and outstanding balances were not always delivered directly.", BLUE),
    ("Weak access and traceability", "Paper files offered limited role control, ownership checks and message history.", GREEN),
]
positions = [(0.72, 1.65), (6.78, 1.65), (0.72, 3.65), (6.78, 3.65)]
for (title, body, accent), (x, y) in zip(cards, positions):
    add_card(slide, title, body, x, y, 5.78, 1.52, accent=accent)
rect(slide, 0.72, 5.73, 11.84, 0.84, GREEN_DARK, GREEN_DARK, True)
textbox(slide, "Need identified", 0.98, 5.94, 1.75, 0.3, size=14, color=GOLD, bold=True)
textbox(slide,
        "An integrated, secure and easy-to-use system that records payments, updates balances, issues receipts, reports to management and communicates directly with parents.",
        2.55, 5.87, 9.58, 0.48, size=14, color=WHITE)
footer(slide, 3)


# Slide 4 — Objectives and scope
slide = prs.slides.add_slide(BLANK)
title_bar(slide, "Objectives and project scope", 4)
rect(slide, 0.68, 1.50, 12.0, 0.94, GREEN_DARK, GREEN_DARK, True)
textbox(slide, "MAIN OBJECTIVE", 0.95, 1.70, 1.55, 0.28, size=12, color=GOLD, bold=True)
textbox(slide,
        "Design, develop and document a web-based fee tracking and reminder system for Mbonea Secondary School.",
        2.55, 1.64, 9.65, 0.45, size=16, color=WHITE, bold=True)
specific = [
    ("Study", "Analyse the existing fee process, users, information and controls."),
    ("Design", "Create a secure database, workflows and role-based system structure."),
    ("Implement", "Deliver registration, fees, payments, receipts, reports, reminders and parent access."),
    ("Validate", "Test the completed workflows and prepare clear technical documentation."),
]
for i, (title, body) in enumerate(specific):
    add_card(slide, title, body, 0.72 + i * 3.04, 2.78, 2.76, 1.55, accent=[BLUE, GREEN, GOLD, RED][i], number=i + 1)
textbox(slide, "Included", 0.75, 4.72, 2.2, 0.3, size=16, color=GREEN_DARK, bold=True)
bullets(slide, [
    "Fee structures, payments, receipts and balances",
    "SMS/email reminders and parent portal",
    "Reports, bank evidence and clearance certificates",
], 0.77, 5.08, 5.5, 1.42, size=14, spacing=7)
textbox(slide, "Outside current scope", 6.85, 4.72, 2.8, 0.3, size=16, color=GREEN_DARK, bold=True)
bullets(slide, [
    "Examinations, attendance and timetables",
    "Full academic / learning management",
    "Public production hosting in this project phase",
], 6.87, 5.08, 5.3, 1.42, size=14, spacing=7)
footer(slide, 4)


# Slide 5 — Stakeholders
slide = prs.slides.add_slide(BLANK)
title_bar(slide, "Stakeholders and user roles", 5, "Role-based access keeps daily operations, configuration and parent records separated")
roles = [
    ("School administrator", "Registers students, records payments, prints receipts, sends reminders and prepares reports.", GREEN),
    ("Super administrator", "Configures fee structures, categories, school settings, templates and bank accounts.", GOLD),
    ("Parent / guardian", "Views linked children, balances, history and notifications; uploads bank evidence.", BLUE),
]
for i, (name, body, accent) in enumerate(roles):
    add_card(slide, name, body, 0.65, 1.55 + i * 1.48, 5.2, 1.14, accent=accent, number=i + 1)
add_picture_contain(slide, ASSETS / "01-login.png", 6.25, 1.48, 6.35, 4.88, True)
label(slide, "Three login paths → one controlled system", 7.62, 6.48, 3.65, GREEN_SOFT, GREEN_DARK, 11)
footer(slide, 5)


# Slide 6 — Methodology
slide = prs.slides.add_slide(BLANK)
title_bar(slide, "Methodology and project plan", 6, "Agile, iterative delivery supported by interviews, observation and document review")
textbox(slide, "Requirements gathering", 0.68, 1.52, 3.2, 0.3, size=17, color=GREEN_DARK, bold=True)
methods = [("Interview", "User roles, reports and communication needs"),
           ("Observation", "Payment and receipt handling steps"),
           ("Document review", "Data fields, reports and receipt formats"),
           ("Group discussion", "Priorities and team task allocation")]
for i, (m, d) in enumerate(methods):
    add_card(slide, m, d, 0.68, 1.94 + i * 0.96, 4.05, 0.78, accent=[GREEN, BLUE, GOLD, RED][i])

textbox(slide, "Six development iterations", 5.18, 1.52, 3.4, 0.3, size=17, color=GREEN_DARK, bold=True)
iterations = [
    "Foundation & roles", "Students & fee records", "Payments & receipts",
    "Communication", "Reports & parent access", "Final testing & documentation"
]
for i, item in enumerate(iterations):
    col, row = i % 2, i // 2
    add_card(slide, item, "", 5.18 + col * 3.58, 1.94 + row * 1.08, 3.30, 0.84,
             accent=GREEN if i % 2 == 0 else GOLD, number=i + 1)

add_metric(slide, "12 weeks", "23 Mar – 15 Jun 2026", 5.18, 5.52, 2.25, GREEN)
add_metric(slide, "TZS 80,000", "Academic project budget", 7.65, 5.52, 2.25, GOLD)
add_metric(slide, "3 members", "Front-end, back-end & shared analysis", 10.12, 5.52, 2.25, BLUE)
footer(slide, 6)


# Slide 7 — Requirements
slide = prs.slides.add_slide(BLANK)
title_bar(slide, "Requirements translated into completed modules", 7)
modules = [
    ("Authentication", "Role-based login and access checks"),
    ("Student & parent", "Register, import and official guardian links"),
    ("Fee setup", "Structures, categories and installment schedule"),
    ("Payments", "Category allocation, balance and numbered receipt"),
    ("Communication", "SMS/email templates, reminders and logs"),
    ("Reports", "Fee position, unpaid, receipt, bank and clearance"),
    ("Parent portal", "Linked children, history, notifications and bank upload"),
    ("Settings", "School identity, message and bank configuration"),
]
for i, (name, body) in enumerate(modules):
    col, row = i % 4, i // 4
    add_card(slide, name, body, 0.63 + col * 3.12, 1.55 + row * 1.72, 2.85, 1.42,
             accent=[GREEN, BLUE, GOLD, RED][i % 4])
rect(slide, 0.66, 5.25, 12.0, 1.18, GREEN_SOFT, GREEN_SOFT, True)
textbox(slide, "NON-FUNCTIONAL REQUIREMENTS", 0.94, 5.48, 2.45, 0.3, size=12, color=GREEN_DARK, bold=True)
textbox(slide,
        "Usability  ·  Performance  ·  Security  ·  Reliability  ·  Browser compatibility  ·  Responsive design",
        3.20, 5.42, 8.92, 0.42, size=15, color=INK, bold=True, align=PP_ALIGN.CENTER)
footer(slide, 7)


# Slide 8 — Architecture
slide = prs.slides.add_slide(BLANK)
title_bar(slide, "System architecture and technology stack", 8, "Laravel MVC with service classes for communication, reporting and bank verification")
layers = [
    ("Web browser", "Responsive Blade + Bootstrap interface", BLUE),
    ("Routes & middleware", "Authentication, CSRF and role checks", GOLD),
    ("Controllers & services", "Validation, fee rules, reports, messages, bank parsing", GREEN),
    ("Models & MySQL", "Related operational records with foreign keys", RED),
]
for i, (name, body, accent) in enumerate(layers):
    y = 1.56 + i * 1.12
    add_card(slide, name, body, 0.72, y, 6.0, 0.86, accent=accent, number=i + 1)
    if i < len(layers) - 1:
        textbox(slide, "↓", 3.45, y + 0.82, 0.4, 0.30, size=18, color=MUTED,
                bold=True, align=PP_ALIGN.CENTER)

rect(slide, 7.20, 1.56, 5.40, 3.43, WHITE, LINE, True)
textbox(slide, "Specialised integrations", 7.52, 1.82, 3.8, 0.34, size=18, color=GREEN_DARK, bold=True)
integration_cards = [
    ("SMS gateway", "Payment confirmations & reminders"),
    ("Email server", "Welcome messages & fee notices"),
    ("DomPDF", "Receipts, reports & clearance"),
    ("Laravel Excel", "Student import & report export"),
]
for i, (name, body) in enumerate(integration_cards):
    col, row = i % 2, i // 2
    add_card(slide, name, body, 7.48 + col * 2.47, 2.34 + row * 1.10, 2.25, 0.86,
             accent=GREEN if i < 2 else GOLD)
textbox(slide, "Implemented technologies", 7.22, 5.28, 3.2, 0.3, size=16, color=GREEN_DARK, bold=True)
techs = ["Laravel 12", "PHP 8.2", "MySQL", "Bootstrap 5", "Blade", "Git/GitHub", "XAMPP"]
for i, tech in enumerate(techs):
    x = 7.22 + (i % 4) * 1.32
    y = 5.72 + (i // 4) * 0.50
    label(slide, tech, x, y, 1.18, GREEN_SOFT, GREEN_DARK, 9.5)
footer(slide, 8)


# Slide 9 — Database/security
slide = prs.slides.add_slide(BLANK)
title_bar(slide, "Database and security design", 9)
add_picture_contain(slide, DIAGRAMS / "word-style" / "png" / "05-er-diagram.png",
                    0.58, 1.43, 7.10, 5.50, True)
textbox(slide, "Security controls", 8.03, 1.53, 3.0, 0.34, size=18, color=GREEN_DARK, bold=True)
security = [
    "Hashed passwords and rate-limited login",
    "Role middleware for admin, super admin and parent",
    "Parent record scope through official student links",
    "Validated requests and CSRF protection",
    "Controlled PDF uploads for bank evidence",
    "Foreign keys and unique values for consistency",
]
bullets(slide, security, 8.04, 2.02, 4.55, 3.25, size=14.5, spacing=9)
rect(slide, 8.02, 5.48, 4.62, 1.04, GREEN_DARK, GREEN_DARK, True)
textbox(slide, "Design principle", 8.25, 5.70, 1.4, 0.28, size=12, color=GOLD, bold=True)
textbox(slide, "Every payment, message and parent link remains traceable to the correct student.",
        9.54, 5.61, 2.74, 0.52, size=12.5, color=WHITE)
footer(slide, 9)


# Slide 10 — Student management
slide = prs.slides.add_slide(BLANK)
title_bar(slide, "Student and parent management flow", 10)
add_picture_crop(slide, ASSETS / "03-students.png", 0.60, 1.43, 8.05, 5.42, True)
textbox(slide, "Implemented workflow", 9.00, 1.53, 3.2, 0.34, size=18, color=GREEN_DARK, bold=True)
bullets(slide, [
    "Register a new student or import from CSV/XLSX.",
    "Select an existing parent account or create a new parent during admission.",
    "Create an official parent–student link with relationship and contact details.",
    "Send a welcome email containing phone login details and a temporary password.",
    "Assign fee structures and apply the school-wide installment schedule.",
], 9.00, 2.02, 3.75, 3.75, size=14, spacing=9)
label(slide, "One parent can be linked to multiple children", 9.18, 6.16, 3.32, GREEN_SOFT, GREEN_DARK, 10.5)
footer(slide, 10)


# Slide 11 — Payment
slide = prs.slides.add_slide(BLANK)
title_bar(slide, "Payment recording and receipt workflow", 11)
add_picture_crop(slide, ASSETS / "04-record-payment.png", 0.58, 1.43, 8.15, 5.42, True)
add_picture_contain(slide, DIAGRAMS / "png" / "02-activity-diagram-fee-lifecycle.png",
                    9.00, 1.43, 3.72, 3.88, True)
bullets(slide, [
    "Validate student, date, mode and category allocations.",
    "Generate a unique term-based receipt number.",
    "Update paid amount and outstanding balance immediately.",
    "Log and send payment confirmation to the parent.",
], 9.02, 5.48, 3.63, 1.18, size=11.5, spacing=4)
footer(slide, 11)


# Slide 12 — Communication
slide = prs.slides.add_slide(BLANK)
title_bar(slide, "SMS, email and reminder workflow", 12)
add_picture_crop(slide, ASSETS / "07-send-reminders.png", 0.58, 1.43, 8.40, 5.42, True)
textbox(slide, "System-driven communication", 9.28, 1.56, 3.25, 0.34, size=18, color=GREEN_DARK, bold=True)
bullets(slide, [
    "New parents appear automatically in the selectable list.",
    "Manual batches are restricted to 1–5 parents.",
    "Templates are ordered by fee status: overdue, due today, 3/7/14 days, general.",
    "Auto mode chooses the correct template for each selected parent.",
    "Scheduled reminders run daily and message history stores the result.",
], 9.28, 2.05, 3.45, 3.45, size=13.5, spacing=8)
rect(slide, 9.27, 5.72, 3.47, 0.80, GREEN_SOFT, GREEN_SOFT, True)
textbox(slide, "Channels", 9.50, 5.93, 0.85, 0.24, size=11, color=GREEN_DARK, bold=True)
textbox(slide, "SMS  +  Email  +  Portal logs", 10.37, 5.87, 2.10, 0.34,
        size=12.5, color=INK, bold=True)
footer(slide, 12)


# Slide 13 — Reports
slide = prs.slides.add_slide(BLANK)
title_bar(slide, "Reports and management information", 13)
add_picture_crop(slide, ASSETS / "08-reports-hub.png", 0.58, 1.43, 8.45, 5.42, True)
textbox(slide, "Live reports from recorded transactions", 9.32, 1.55, 3.18, 0.60,
        size=18, color=GREEN_DARK, bold=True)
report_items = [
    "School fee position",
    "Collection by period",
    "Receipt register",
    "Unpaid balances",
    "Paid / clearance",
    "Message history",
    "Bank proof review",
]
bullets(slide, report_items, 9.34, 2.30, 3.1, 2.95, size=14, spacing=7)
label(slide, "Browser view", 9.34, 5.45, 1.18, GREEN_SOFT, GREEN_DARK, 9.5)
label(slide, "PDF", 10.68, 5.45, 0.72, RGBColor(250, 231, 232), RED, 9.5)
label(slide, "Excel", 11.57, 5.45, 0.85, GREEN_SOFT, GREEN_DARK, 9.5)
rect(slide, 9.33, 6.05, 3.10, 0.55, GREEN_DARK, GREEN_DARK, True)
textbox(slide, "Supports bursar follow-up and management oversight",
        9.49, 6.16, 2.78, 0.26, size=10.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
footer(slide, 13)


# Slide 14 — Parent portal
slide = prs.slides.add_slide(BLANK)
title_bar(slide, "Parent portal and bank payment workflow", 14)
add_picture_crop(slide, ASSETS / "12-parent-dashboard.png", 0.55, 1.40, 7.38, 3.32, True)
add_picture_crop(slide, ASSETS / "14-parent-bank-payment.png", 0.55, 4.91, 7.38, 1.89, True)
add_picture_contain(slide, ASSETS / "15-parent-mobile.png", 8.20, 1.40, 2.75, 5.40, True)
textbox(slide, "Parent capabilities", 11.18, 1.54, 1.72, 0.34, size=17, color=GREEN_DARK, bold=True)
bullets(slide, [
    "View only officially linked children",
    "See expected, paid and outstanding amounts",
    "Review receipts and notification history",
    "Upload NMB/CRDB receipt PDFs",
    "Download clearance when fully paid",
    "Use desktop or mobile interface",
], 11.15, 2.06, 1.75, 3.62, size=11.5, spacing=8)
label(slide, "Phone + password login", 11.17, 6.16, 1.67, GREEN_SOFT, GREEN_DARK, 9.5)
footer(slide, 14)


# Slide 15 — Testing/conclusion
slide = prs.slides.add_slide(BLANK)
title_bar(slide, "Testing, conclusion and next steps", 15)
add_metric(slide, "104", "Automated tests passed", 0.70, 1.48, 2.25, GREEN)
add_metric(slide, "329", "Assertions passed", 3.12, 1.48, 2.25, BLUE)
add_metric(slide, "12", "Report test areas passed", 5.54, 1.48, 2.25, GOLD)
add_metric(slide, "3 roles", "Access boundaries verified", 7.96, 1.48, 2.25, RED)
add_metric(slide, "1–5", "Parent batch rule verified", 10.38, 1.48, 2.25, GREEN)

textbox(slide, "Conclusion", 0.72, 2.94, 2.0, 0.35, size=18, color=GREEN_DARK, bold=True)
bullets(slide, [
    "The completed system centralises fee records, balances, receipts, reports and parent communication.",
    "Role controls and official parent links protect financial and student information.",
    "All stated project objectives and report test areas were achieved for local demonstration.",
], 0.74, 3.37, 5.62, 2.25, size=14, spacing=9)

textbox(slide, "Limitations and recommendations", 6.82, 2.94, 3.9, 0.35, size=18, color=GREEN_DARK, bold=True)
bullets(slide, [
    "Deploy to a secure production server with automated backups.",
    "Configure live SMS/email credentials and monitor delivery.",
    "Improve bank integrations as structured bank data becomes available.",
    "Provide staff training and periodic security review.",
], 6.84, 3.37, 5.55, 2.25, size=14, spacing=9)

rect(slide, 0.70, 6.10, 11.93, 0.58, GREEN_DARK, GREEN_DARK, True)
textbox(slide, "THANK YOU  ·  QUESTIONS & DEMONSTRATION",
        0.95, 6.22, 11.42, 0.28, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
footer(slide, 15, "Sources: Final Year Project Report, implemented system screenshots & automated test suite · July 2026")


output = OUT_DIR / "Web-Based-Fee-Tracking-and-Reminder-System-Presentation.pptx"
prs.save(output)
print(f"Saved {output}")
print(f"Slides: {len(prs.slides)}")
